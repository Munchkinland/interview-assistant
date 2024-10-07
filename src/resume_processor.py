import os
import docx2txt
import PyPDF2
from pptx import Presentation

class ResumeProcessor:
  @staticmethod
  def process_resume(file_path):
      _, file_extension = os.path.splitext(file_path)
      
      if file_extension.lower() == '.txt':
          with open(file_path, 'r', encoding='utf-8') as file:
              return file.read()
      elif file_extension.lower() == '.docx':
          return docx2txt.process(file_path)
      elif file_extension.lower() == '.pdf':
          with open(file_path, 'rb') as file:
              pdf_reader = PyPDF2.PdfReader(file)
              return '\n'.join([page.extract_text() for page in pdf_reader.pages])
      elif file_extension.lower() in ['.ppt', '.pptx']:
          prs = Presentation(file_path)
          text = []
          for slide in prs.slides:
              for shape in slide.shapes:
                  if hasattr(shape, 'text'):
                      text.append(shape.text)
          return '\n'.join(text)
      else:
          raise ValueError(f"Formato de archivo no soportado: {file_extension}")